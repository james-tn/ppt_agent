###############################################################################  
# app.py  –  Chat-driven AI slide-deck generator that outputs PowerPoint (.pptx)  
###############################################################################  
import os, base64, mimetypes, json, textwrap, re, io  
from pathlib import Path  
  
import streamlit as st  
import streamlit.components.v1 as components  
from dotenv import load_dotenv  
from openai import AzureOpenAI                # pip install openai>=1.12.0  
from pptx import Presentation                 # pip install python-pptx  
from pptx.util import Inches, Pt  
  
###############################################################################  
# 1.  Azure OpenAI client  
###############################################################################  
load_dotenv()  
client = AzureOpenAI(  
    api_key        = os.getenv("AZURE_OPENAI_API_KEY"),  
    api_version    = os.getenv("AZURE_OPENAI_API_VERSION"),  
    azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT"),  
)  
MODEL_NAME = os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT")  
  
###############################################################################  
# 2.  Helper functions  
###############################################################################  
IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}  
  
def image_to_b64(p: Path) -> str:  
    return base64.b64encode(p.read_bytes()).decode("utf-8")  
  
def read_folder(folder: Path):  
    """Return (markdown_text, images_dict {filename: b64}) for every file."""  
    md_chunks, image_dict = [], {}  
    for fp in folder.rglob("*"):  
        if fp.suffix.lower() == ".md":  
            md_chunks.append(fp.read_text(encoding="utf-8", errors="ignore"))  
        elif fp.suffix.lower() in IMG_EXT:  
            image_dict[fp.name] = image_to_b64(fp)  
    return "\n\n".join(md_chunks), image_dict  
  
def build_initial_messages(prompt: str, md: str, images: dict):  
    """First call – supply markdown text + pictures as context."""  
    sys_msg = {  
        "role": "system",  
        "content": textwrap.dedent("""  
            You are SlideBuilder-GPT.  
            Return ONLY valid JSON (no markdown, no code fences).  
            JSON schema:  
              {  
                "slides":[  
                  {  
                    "title":"Slide title",  
                    "points":["bullet 1","bullet 2"],  
                    "images":["file1.png","file2.jpg"],        // optional  
                    "notes":"optional speaker notes"  
                  }  
                ]  
              }  
            Rules:  
            • The file names in the images array MUST match exactly the files  
              you have been shown (case-sensitive).  
            • Do NOT embed base-64.  Do NOT output HTML.  
        """).strip(),  
    }  
  
    user_parts = [{"type": "text", "text": prompt}]  
    if md.strip():  
        user_parts.append({"type": "text",  
                           "text": f"\n\nHere are the markdown files:\n\n{md}"})  
    for fn, b64 in images.items():  
        mime, _ = mimetypes.guess_type(fn)  
        mime = mime or "image/png"  
        # show model the actual picture  
        user_parts.append({  
            "type": "image_url",  
            "image_url": {"url": f"data:{mime};base64,{b64}"},  
        })  
        user_parts.append({"type": "text", "text": f"(filename: {fn})"})  
  
    return [sys_msg, {"role": "user", "content": user_parts}]  
  
def openai_call(messages):  
    resp = client.chat.completions.create(  
        model       = MODEL_NAME,  
        messages    = messages,  
    )  
    return resp.choices[0].message.content.strip()  
  
def parse_slides_json(raw: str) -> dict:  
    """  
    Ensure the assistant reply is valid JSON.  
    If the model wraps it in ```json``` fences, strip them.  
    """  
    raw = raw.strip()  
    if raw.startswith("```"):  
        raw = re.sub(r"^```[a-zA-Z]*", "", raw).rstrip("`").strip()  
    return json.loads(raw)  
  
# -------------------------- PowerPoint generation -------------------------  
  
def build_pptx(slides_dict: dict, folder: Path) -> bytes:  
    prs = Presentation()  
    title_content_layout = prs.slide_layouts[1]   # Title and Content  
  
    for slide_data in slides_dict.get("slides", []):  
        title_text   = slide_data.get("title", "")  
        bullets      = slide_data.get("points", []) or []  
        images       = slide_data.get("images", []) or []  
        notes_text   = slide_data.get("notes", "")  
  
        slide = prs.slides.add_slide(title_content_layout)  
  
        # Title  
        slide.shapes.title.text = title_text  
  
        # Bullets  
        body = slide.shapes.placeholders[1].text_frame  
        body.clear()                 # remove default bullet  
        for idx, b in enumerate(bullets):  
            p = body.add_paragraph() if idx else body.paragraphs[0]  
            p.text  = b  
            p.level = 0  
            p.font.size = Pt(18)  
  
        # Images – stack under the body placeholder  
        pic_left = Inches(5.5)      # right-hand side  
        pic_top  = Inches(1.5)  
        maxw     = Inches(3.0)  
        for img in images:  
            img_path = folder / img  
            if img_path.exists():  
                slide.shapes.add_picture(str(img_path), pic_left, pic_top,  
                                          width=maxw)  
                pic_top += Inches(2.5)  
  
        # Speaker notes  
        if notes_text:  
            notes = slide.notes_slide.notes_text_frame  
            notes.text = notes_text  
  
    bio = io.BytesIO()  
    prs.save(bio)  
    bio.seek(0)  
    return bio.read()  
  
def json_to_html_preview(slides_dict: dict) -> str:  
    """Crude, fast preview: turn JSON into a very simple HTML list."""  
    html = ["<style>h2{margin:4px;}ul{margin-top:0}</style>"]  
    for s in slides_dict.get("slides", []):  
        html.append(f"<h2>{s.get('title','(no title)')}</h2>")  
        if s.get("points"):  
            html.append("<ul>" + "".join(f"<li>{p}</li>" for p in s["points"])  
                        + "</ul>")  
        for img in s.get("images", []):  
            html.append(f'<div style="margin-bottom:6px;">🖼 {img}</div>')  
        html.append("<hr>")  
    return "".join(html)  
  
###############################################################################  
# 3.  Streamlit UI / State  
###############################################################################  
st.set_page_config(page_title="AI PPTX Generator", layout="wide")  
  
for k, d in [("messages", []), ("openai_messages", []),  
             ("pptx_bytes", b""), ("slides_json", {}),  
             ("folder", "")]:  
    st.session_state.setdefault(k, d)  
  
# ----------------------------  sidebar  ------------------------------------  
with st.sidebar:  
    st.header("📁 Source folder")  
    folder_input = st.text_input("Path with .md + images",  
                                 value=st.session_state["folder"])  
    if folder_input != st.session_state["folder"]:  
        st.session_state["folder"]        = folder_input  
        st.session_state["messages"]      = []  
        st.session_state["openai_messages"] = []  
        st.session_state["pptx_bytes"]    = b""  
        st.session_state["slides_json"]   = {}  
  
    if st.button("🔄 Reset conversation"):  
        for k in ["messages", "openai_messages", "pptx_bytes", "slides_json"]:  
            st.session_state[k] = [] if "messages" in k else b""  
        st.experimental_rerun()  
  
st.title("🧑‍🏫  AI Slide-Deck Generator  →  PowerPoint (.pptx)")  
  
# ----------------- show chat history ---------------------------------------  
for m in st.session_state["messages"]:  
    with st.chat_message(m["role"]):  
        st.markdown(m["content"] if m["role"] == "user"  
                    else "Deck updated ↓")  
  
# ----------------- chat input (generate / refine) --------------------------  
user_msg = st.chat_input("Type the prompt (first time) or feedback (later)…")  
  
if user_msg:  
    # show immediately  
    st.session_state["messages"].append({"role": "user", "content": user_msg})  
    with st.chat_message("user"):  
        st.markdown(user_msg)  
  
    folder_path = Path(st.session_state["folder"]).expanduser()  
    if not folder_path.is_dir():  
        st.error("Please enter a valid folder path in the sidebar before chat.")  
        st.stop()  
  
    # first generation or refinement?  
    first_time = not st.session_state["pptx_bytes"]  
  
    if first_time:  
        md_text, images_dict = read_folder(folder_path)  
        st.session_state["openai_messages"] = build_initial_messages(  
            user_msg, md_text, images_dict)  
    else:  
        st.session_state["openai_messages"].append(  
            {"role": "user", "content": user_msg})  
  
    # call OpenAI  
    with st.spinner("Building deck with GPT…"):  
        raw_reply = openai_call(st.session_state["openai_messages"])  
  
    # append assistant raw JSON to conversation  
    st.session_state["openai_messages"].append(  
        {"role": "assistant", "content": raw_reply})  
  
    # convert JSON → pptx  
    try:  
        slides_dict = parse_slides_json(raw_reply)  
        pptx_data   = build_pptx(slides_dict, folder_path)  
    except Exception as ex:  
        st.error(f"Could not build PowerPoint: {ex}")  
        st.stop()  
  
    st.session_state["slides_json"]  = slides_dict  
    st.session_state["pptx_bytes"]   = pptx_data  
    st.session_state["messages"].append(  
        {"role": "assistant", "content": "Deck updated ↓"})  
  
    with st.chat_message("assistant"):  
        st.markdown("Here’s the updated deck ↓")  
  
# ----------------- preview & download --------------------------------------  
if st.session_state["pptx_bytes"]:  
    st.subheader("🔍 Quick preview")  
    components.html(  
        json_to_html_preview(st.session_state["slides_json"]),  
        height=500,  
        scrolling=True,  
    )  
  
    st.download_button(  
        label="💾 Download .pptx",  
        data=st.session_state["pptx_bytes"],  
        file_name="deck.pptx",  
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",  
    )  
else:  
    st.info("Send the first prompt in the chat box to create your deck.")  